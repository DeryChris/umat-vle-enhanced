# ============================================================
# Multi-format document loader: PDF, DOCX, PPTX, XLSX, TXT, code, audio/video (via Whisper)
# ============================================================

from langchain.text_splitter import RecursiveCharacterTextSplitter
from pathlib import Path
from typing import List, Tuple, Optional
from config import get_settings
import uuid

settings = get_settings()

# Lazily-imported heavy libraries
_pypdf = None
_docx  = None
_pptx  = None
_xl    = None
_whisper_model = None


def _import_pypdf():
    global _pypdf
    if _pypdf is None:
        from pypdf import PdfReader as R
        _pypdf = R
    return _pypdf


def _import_docx():
    global _docx
    if _docx is None:
        from docx import Document as D
        _docx = D
    return _docx


def _import_pptx():
    global _pptx
    if _pptx is None:
        from pptx import Presentation as P
        _pptx = P
    return _pptx


def _import_openpyxl():
    global _xl
    if _xl is None:
        from openpyxl import load_workbook as L
        _xl = L
    return _xl


def _get_whisper():
    global _whisper_model
    if _whisper_model is None:
        import whisper
        _whisper_model = whisper.load_model(settings.whisper_model)
    return _whisper_model


CODE_EXTENSIONS = {
    '.py', '.js', '.ts', '.jsx', '.tsx', '.php', '.rb', '.go', '.rs',
    '.java', '.kt', '.scala', '.swift', '.c', '.cpp', '.h', '.hpp',
    '.cs', '.sql', '.sh', '.bash', '.zsh', '.ps1', '.bat', '.pl',
    '.pm', '.lua', '.r', '.m', '.mm', '.dart', '.html', '.htm',
    '.css', '.scss', '.less', '.sass', '.json', '.xml', '.yaml',
    '.yml', '.toml', '.ini', '.cfg', '.conf', '.md', '.rst', '.tex',
    '.dockerfile', '.makefile', '.cmake', '.gradle',
}

AUDIO_VIDEO_EXTENSIONS = {'.mp3', '.wav', '.ogg', '.flac', '.m4a',
                          '.mp4', '.webm', '.mov', '.avi', '.mkv'}


class DocumentLoader:

    def __init__(self):
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings.max_chunk_size,
            chunk_overlap=settings.chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""],
            length_function=len,
        )

    # ── Individual format loaders ──────────────────────────

    def load_pdf(self, file_path: str) -> str:
        PdfReader = _import_pypdf()
        reader = PdfReader(file_path)
        parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text:
                parts.append(f"[Page {i + 1}]\n{text}")
        return "\n\n".join(parts)

    def load_docx(self, file_path: str) -> str:
        Document = _import_docx()
        doc = Document(file_path)
        parts = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                parts.append(t)
        # Include table content
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))
        return "\n\n".join(parts)

    def load_pptx(self, file_path: str) -> str:
        Presentation = _import_pptx()
        filename = Path(file_path).name
        parts = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_texts = [f"--- Slide {i + 1} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_texts.append(" | ".join(cells))
                parts.append("\n".join(slide_texts))
        except Exception:
            return f"[PowerPoint: {filename}]"

        if not parts:
            return f"[PowerPoint: {filename}]"

        return "\n\n".join(parts)

    def load_old_doc(self, file_path: str) -> str:
        """Extract text from old .doc (Word 97-2003) via LibreOffice or basic OleFile."""
        import subprocess, tempfile, logging
        logger = logging.getLogger(__name__)
        try:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'txt:Text', '--outdir', tempfile.gettempdir(), file_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                out = Path(tempfile.gettempdir()) / (Path(file_path).stem + '.txt')
                if out.exists():
                    text = out.read_text(encoding='utf-8', errors='ignore')
                    out.unlink(missing_ok=True)
                    if text.strip():
                        return text
        except Exception:
            pass
        logger.warning(f"LibreOffice not available for {file_path}; trying OleFile fallback")
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)
            text = ''
            for stream_name in ole.listdir():
                name = '/'.join(stream_name)
                try:
                    data = ole.openstream(name).read()
                    decoded = data.decode('utf-16', errors='ignore')
                    text += decoded + '\n'
                except Exception:
                    try:
                        decoded = data.decode('utf-8', errors='ignore')
                        text += decoded + '\n'
                    except Exception:
                        pass
            ole.close()
            clean = '\n'.join(line for line in text.split('\n') if line.strip())
            if clean.strip():
                return clean
        except ImportError:
            logger.warning("olefile not installed. Install with: pip install olefile")
        except Exception:
            pass
        return f"[Could not extract text from {Path(file_path).name}. Convert to .docx format for better results.]"

    def load_old_ppt(self, file_path: str) -> str:
        """Extract text from old .ppt (PowerPoint 97-2003) via LibreOffice or basic OleFile."""
        import subprocess, tempfile, logging
        logger = logging.getLogger(__name__)
        try:
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'txt:Text', '--outdir', tempfile.gettempdir(), file_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                out = Path(tempfile.gettempdir()) / (Path(file_path).stem + '.txt')
                if out.exists():
                    text = out.read_text(encoding='utf-8', errors='ignore')
                    out.unlink(missing_ok=True)
                    if text.strip():
                        return text
        except Exception:
            pass
        logger.warning(f"LibreOffice not available for {file_path}; trying OleFile fallback")
        try:
            import olefile
            ole = olefile.OleFileIO(file_path)
            text = ''
            for stream_name in ole.listdir():
                name = '/'.join(stream_name)
                try:
                    data = ole.openstream(name).read()
                    decoded = data.decode('utf-16', errors='ignore')
                    text += decoded + '\n'
                except Exception:
                    try:
                        decoded = data.decode('utf-8', errors='ignore')
                        text += decoded + '\n'
                    except Exception:
                        pass
            ole.close()
            clean = '\n'.join(line for line in text.split('\n') if line.strip())
            if clean.strip():
                return clean
        except ImportError:
            logger.warning("olefile not installed. Install with: pip install olefile")
        except Exception:
            pass
        return f"[Could not extract text from {Path(file_path).name}. Convert to .pptx format for better results.]"

    def load_xlsx(self, file_path: str) -> str:
        load_workbook = _import_openpyxl()
        wb = load_workbook(file_path, read_only=True, data_only=True)
        parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = [f"=== Sheet: {sheet_name} ==="]
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows_text.append("\t".join(cells))
            parts.append("\n".join(rows_text))
        wb.close()
        return "\n\n".join(parts)

    def load_text(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    def load_code(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        lang = ext.lstrip(".") if ext else "txt"
        text = self.load_text(file_path)
        return f"[File: {Path(file_path).name}]\n[Language: {lang}]\n\n{text}"

    def load_media(self, file_path: str) -> str:
        """Transcribe audio/video via Whisper and return text."""
        model = _get_whisper()
        result = model.transcribe(file_path, language="en")
        text = result.get("text", "")
        segments = result.get("segments", [])
        if segments:
            lines = []
            for seg in segments:
                start = seg.get("start", 0)
                end = seg.get("end", 0)
                txt = seg.get("text", "").strip()
                if txt:
                    lines.append(f"[{_fmt_ts(start)} - {_fmt_ts(end)}] {txt}")
            text = "\n".join(lines)
        return f"[Media transcription: {Path(file_path).name}]\n\n{text}"

    # ── Public entry point ────────────────────────────────

    def load_file(self, file_path: str, ext_hint: Optional[str] = None) -> str:
        """
        Detect file type and extract text content.

        Args:
            file_path: Path to the file on disk.
            ext_hint:  Optional extension override (e.g. if original had no extension).
                       Otherwise inferred from the filename suffix.

        Returns:
            Extracted text content as a single string.
        """
        ext = (ext_hint or Path(file_path).suffix).lower()

        if ext == ".pdf":
            return self.load_pdf(file_path)
        elif ext == ".docx":
            return self.load_docx(file_path)
        elif ext == ".doc":
            return self.load_old_doc(file_path)
        elif ext == ".pptx":
            return self.load_pptx(file_path)
        elif ext == ".ppt":
            return self.load_old_ppt(file_path)
        elif ext == ".xlsx":
            return self.load_xlsx(file_path)
        elif ext in {".txt", ".text", ".md", ".markdown", ".csv"}:
            return self.load_text(file_path)
        elif ext in CODE_EXTENSIONS:
            return self.load_code(file_path)
        elif ext in AUDIO_VIDEO_EXTENSIONS:
            return self.load_media(file_path)
        else:
            # Fallback: try as text
            return self.load_text(file_path)

    # ── Chunking ──────────────────────────────────────────

    def split_text(
        self,
        text:     str,
        metadata: dict,
    ) -> Tuple[List[str], List[dict], List[str]]:
        chunks = self.splitter.split_text(text)
        texts, metadatas, ids = [], [], []

        for i, chunk in enumerate(chunks):
            texts.append(chunk)
            metadatas.append({
                **metadata,
                "chunk_index":  i,
                "total_chunks": len(chunks),
            })
            ids.append(f"{metadata.get('source', 'doc')}_{i}_{uuid.uuid4().hex[:8]}")

        return texts, metadatas, ids

    def process_transcript(
        self,
        transcript: str,
        session_id: str,
        course_id:  int,
    ) -> Tuple[List[str], List[dict], List[str]]:
        metadata = {
            "source":      f"transcript_{session_id}",
            "source_type": "transcript",
            "session_id":  session_id,
            "course_id":   str(course_id),
        }
        return self.split_text(transcript, metadata)


def _fmt_ts(seconds: float) -> str:
    m, s = divmod(int(seconds), 60)
    h, m = divmod(m, 60)
    if h:
        return f"{h}:{m:02d}:{s:02d}"
    return f"{m}:{s:02d}"
